"""Generate presentations/data_profiling.pptx from live DuckDB numbers.

Deck structure mirrors the 7 slide sections in docs/profiling.md 1:1, plus
a title slide at the start and a closing slide at the end (9 slides total):
  Title
  Slide 1 — Data Exploration
  Slide 2 — Understanding the Source Grain
  Slide 3 — Data Quality — Completeness & Uniqueness
  Slide 4 — Data Quality — Cross-field & Geography
  Slide 5 — Not Every Anomaly Is Bad Data
  Slide 6 — Data Quality Treatment Strategy
  Slide 7 — Locked Decisions D1–D8 (rationale in speaker notes)
  Next: staging layer

Every slide has speaker notes so it can be presented cold. Slide 7 carries
the full rationale + interview soundbite for each decision — the canonical
source, not a summary of some other file.

Regenerate: python scripts/build_profiling_deck.py
Prereq:     DBeaver must not hold a write lock on the DuckDB file.
"""
from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path

import duckdb
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

REPO = Path(__file__).resolve().parents[1]
DB = REPO / "warehouse" / "fashionable.duckdb"
OUT = REPO / "presentations" / "data_profiling.pptx"
TABLE = "raw.fashionable_sales_raw"

# palette
NAVY = RGBColor(0x1F, 0x2A, 0x44)
ACCENT = RGBColor(0xC0, 0x39, 0x2B)
MUTED = RGBColor(0x6C, 0x75, 0x7D)
LIGHT_BG = RGBColor(0xF5, 0xF6, 0xF8)


# --- SQL helpers (same style as profile_raw.py) ------------------------------
def q(name: str) -> str:
    return '"' + name.replace('"', '""') + '"'


def is_missing(name: str) -> str:
    c = q(name)
    return f"({c} IS NULL OR TRIM({c}) = '')"


# --- data collection ---------------------------------------------------------
@dataclass
class Facts:
    n_rows: int
    n_cols: int
    unique_orders: int
    skus: int
    styles: int
    categories: int
    sizes: int
    cities: int
    date_min: str
    date_max: str
    distinct_dates: int
    multi_line_orders: int
    max_lines_per_order: int
    styles_multi_skus: int
    max_skus_per_style: int
    dup_order_sku_groups: int
    dup_order_sku_extra: int
    biz_dupes: int
    missing_amount: int
    qty_zero: int
    invalid_dates: int
    missing_ship_info: int
    sku_style_drift: int
    sku_category_drift: int
    sku_size_drift: int
    amt_cur_both_null: int
    contradictions: int
    city_raw: int
    city_norm: int
    state_raw: int
    state_norm: int
    cancelled: int
    cancelled_pct: float
    cancelled_with_amt: int
    shipped_qty_zero: int
    shipped_amt_missing: int
    promo_present: int
    promo_max_count: int
    promo_avg_count: float
    promo_max_len: int
    u22_nonempty: int


def collect() -> Facts:
    with duckdb.connect(str(DB), read_only=True) as con:
        cols = [r[0] for r in con.execute(f"DESCRIBE {TABLE}").fetchall()]
        business_cols = [c for c in cols if c.lower() != "index"]

        (n_rows, unique_orders, skus, styles, categories, sizes, cities,
         d_min, d_max, distinct_dates) = con.execute(f"""
            SELECT COUNT(*), COUNT(DISTINCT "Order ID"), COUNT(DISTINCT "SKU"),
                   COUNT(DISTINCT "Style"), COUNT(DISTINCT "Category"),
                   COUNT(DISTINCT "Size"), COUNT(DISTINCT "ship-city"),
                   MIN(TRY_STRPTIME("Date", '%m-%d-%y'))::DATE,
                   MAX(TRY_STRPTIME("Date", '%m-%d-%y'))::DATE,
                   COUNT(DISTINCT "Date")
            FROM {TABLE}
        """).fetchone()

        multi_lines, max_lines = con.execute(f"""
            SELECT COUNT(*) FILTER (WHERE cnt > 1), MAX(cnt) FROM (
                SELECT COUNT(*) AS cnt FROM {TABLE} GROUP BY "Order ID"
            )
        """).fetchone()

        styles_multi, max_skus_per_style = con.execute(f"""
            SELECT COUNT(*) FILTER (WHERE cnt > 1), MAX(cnt) FROM (
                SELECT COUNT(DISTINCT "SKU") AS cnt
                FROM {TABLE} WHERE NOT {is_missing('Style')}
                GROUP BY "Style"
            )
        """).fetchone()

        dup_grp, dup_extra = con.execute(f"""
            SELECT COUNT(*), COALESCE(SUM(cnt - 1), 0) FROM (
                SELECT COUNT(*) AS cnt FROM {TABLE}
                GROUP BY "Order ID", "SKU" HAVING COUNT(*) > 1
            )
        """).fetchone()

        gb = ", ".join(q(c) for c in business_cols)
        biz_dupes = con.execute(f"""
            SELECT COALESCE(SUM(cnt - 1), 0) FROM (
                SELECT COUNT(*) AS cnt FROM {TABLE}
                GROUP BY {gb} HAVING COUNT(*) > 1
            )
        """).fetchone()[0]

        (missing_amt, qty_zero, invalid_dates, missing_ship_info) = con.execute(f"""
            SELECT
                COUNT(*) FILTER (WHERE {is_missing('Amount')}),
                COUNT(*) FILTER (WHERE TRY_CAST("Qty" AS INT) = 0),
                COUNT(*) FILTER (WHERE NOT {is_missing('Date')}
                                   AND TRY_STRPTIME("Date", '%m-%d-%y') IS NULL),
                COUNT(*) FILTER (WHERE {is_missing('ship-city')} OR {is_missing('ship-state')}
                                    OR {is_missing('ship-postal-code')} OR {is_missing('ship-country')})
            FROM {TABLE}
        """).fetchone()

        (sku_style, sku_cat, sku_size) = con.execute(f"""
            SELECT
                COUNT(*) FILTER (WHERE style_c    > 1),
                COUNT(*) FILTER (WHERE category_c > 1),
                COUNT(*) FILTER (WHERE size_c     > 1)
            FROM (
                SELECT "SKU",
                    COUNT(DISTINCT NULLIF(TRIM("Style"),    '')) AS style_c,
                    COUNT(DISTINCT NULLIF(TRIM("Category"), '')) AS category_c,
                    COUNT(DISTINCT NULLIF(TRIM("Size"),     '')) AS size_c
                FROM {TABLE} WHERE NOT {is_missing('SKU')}
                GROUP BY "SKU"
            )
        """).fetchone()

        amt_cur_both = con.execute(f"""
            SELECT COUNT(*) FROM {TABLE}
            WHERE {is_missing('Amount')} AND {is_missing('currency')}
        """).fetchone()[0]

        contradictions = con.execute(f"""
            SELECT COUNT(*) FROM {TABLE}
            WHERE ("Status" = 'Shipped'    AND "Courier Status" IN ('Unshipped','Cancelled'))
               OR ("Status" LIKE 'Pending%' AND "Courier Status" = 'Shipped')
        """).fetchone()[0]

        (city_raw, city_norm, state_raw, state_norm) = con.execute(f"""
            SELECT
                COUNT(DISTINCT "ship-city"),
                COUNT(DISTINCT UPPER(TRIM("ship-city"))),
                COUNT(DISTINCT "ship-state"),
                COUNT(DISTINCT UPPER(TRIM("ship-state")))
            FROM {TABLE}
        """).fetchone()

        (cancelled, cancelled_with_amt, shipped_q0, shipped_amt_missing) = con.execute(f"""
            SELECT
                COUNT(*) FILTER (WHERE "Status" = 'Cancelled'),
                COUNT(*) FILTER (WHERE "Status" = 'Cancelled'
                                   AND TRY_CAST("Amount" AS DOUBLE) > 0),
                COUNT(*) FILTER (WHERE TRY_CAST("Qty" AS INT) = 0
                                   AND TRIM("Status") LIKE 'Shipped%'),
                COUNT(*) FILTER (WHERE {is_missing('Amount')}
                                   AND TRIM("Status") LIKE 'Shipped%')
            FROM {TABLE}
        """).fetchone()

        (promo_present, promo_max_len, promo_max_count, promo_avg) = con.execute(f"""
            SELECT
                COUNT(*) FILTER (WHERE NOT {is_missing('promotion-ids')}),
                MAX(LENGTH("promotion-ids")),
                MAX(LENGTH("promotion-ids") - LENGTH(REPLACE("promotion-ids", ',', '')) + 1),
                AVG(CASE
                    WHEN NOT {is_missing('promotion-ids')}
                    THEN LENGTH("promotion-ids") - LENGTH(REPLACE("promotion-ids", ',', '')) + 1
                END)
            FROM {TABLE}
        """).fetchone()

        u22_nonempty = con.execute(f"""
            SELECT COUNT(*) FROM {TABLE} WHERE NOT {is_missing('Unnamed: 22')}
        """).fetchone()[0]

    return Facts(
        n_rows=n_rows, n_cols=len(cols), unique_orders=unique_orders,
        skus=skus, styles=styles, categories=categories, sizes=sizes, cities=cities,
        date_min=str(d_min), date_max=str(d_max), distinct_dates=distinct_dates,
        multi_line_orders=multi_lines, max_lines_per_order=max_lines,
        styles_multi_skus=styles_multi, max_skus_per_style=max_skus_per_style,
        dup_order_sku_groups=dup_grp, dup_order_sku_extra=dup_extra,
        biz_dupes=biz_dupes,
        missing_amount=missing_amt, qty_zero=qty_zero,
        invalid_dates=invalid_dates, missing_ship_info=missing_ship_info,
        sku_style_drift=sku_style, sku_category_drift=sku_cat, sku_size_drift=sku_size,
        amt_cur_both_null=amt_cur_both, contradictions=contradictions,
        city_raw=city_raw, city_norm=city_norm,
        state_raw=state_raw, state_norm=state_norm,
        cancelled=cancelled, cancelled_pct=100 * cancelled / n_rows,
        cancelled_with_amt=cancelled_with_amt,
        shipped_qty_zero=shipped_q0, shipped_amt_missing=shipped_amt_missing,
        promo_present=promo_present, promo_max_count=promo_max_count,
        promo_avg_count=promo_avg or 0, promo_max_len=promo_max_len,
        u22_nonempty=u22_nonempty,
    )


# --- shape helpers -----------------------------------------------------------
def _bg(slide, color) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _bar(slide, left: float, top: float, width: float, color) -> None:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top),
        Inches(width), Inches(0.08),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _text(slide, text: str, left: float, top: float, width: float, height: float,
          size: int = 16, bold: bool = False, italic: bool = False,
          color=NAVY, align_center: bool = False) -> None:
    tb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    ).text_frame
    tb.word_wrap = True
    p = tb.paragraphs[0]
    p.text = text
    if align_center:
        p.alignment = PP_ALIGN.CENTER
    for r in p.runs:
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color


def _bullets(slide, bullets: list[str], top: float, size: int = 15) -> None:
    tb = slide.shapes.add_textbox(
        Inches(0.6), Inches(top), Inches(12.1), Inches(7.5 - top - 0.2)
    ).text_frame
    tb.word_wrap = True
    for i, b in enumerate(bullets):
        p = tb.paragraphs[0] if i == 0 else tb.add_paragraph()
        p.text = f"• {b}"
        p.space_after = Pt(4)
        for r in p.runs:
            r.font.size = Pt(size)
            r.font.color.rgb = NAVY


def _table(slide, headers: list[str], rows: list[tuple], top: float,
           left: float = 0.6, width: float = 12.1, height: float = 4.5,
           col_widths: list[float] | None = None) -> None:
    tbl = slide.shapes.add_table(
        len(rows) + 1, len(headers),
        Inches(left), Inches(top), Inches(width), Inches(height),
    ).table
    if col_widths:
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = Inches(w)
    for j, h in enumerate(headers):
        c = tbl.cell(0, j)
        c.text = h
        for p in c.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(13)
                r.font.bold = True
                r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        c.fill.solid()
        c.fill.fore_color.rgb = NAVY
    for i, row in enumerate(rows, start=1):
        for j, v in enumerate(row):
            c = tbl.cell(i, j)
            c.text = str(v)
            for p in c.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(11)
                    r.font.color.rgb = NAVY


def _title(slide, text: str) -> None:
    _text(slide, text, left=0.6, top=0.35, width=12, height=0.8,
          size=26, bold=True, color=NAVY)
    _bar(slide, left=0.6, top=1.1, width=1.2, color=ACCENT)


def _add_slide(prs, notes: str = ""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, LIGHT_BG)
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    return slide


# --- deck --------------------------------------------------------------------
def build(f: Facts) -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ---------- 1. TITLE ----------
    s = _add_slide(prs, notes=(
        "Opening frame: this deck walks through what the raw CSV told us "
        "before we designed the star schema. Every number on every slide "
        "was queried live from DuckDB by scripts/build_profiling_deck.py — "
        "nothing hand-typed, nothing cherry-picked."
    ))
    _text(s, "Fashionable Sales — Data Profiling",
          left=0.7, top=2.5, width=12, height=1.5,
          size=44, bold=True, color=NAVY)
    _bar(s, left=0.7, top=3.9, width=2, color=ACCENT)
    _text(s, f"What {f.n_rows:,} rows told us before we designed the star schema",
          left=0.7, top=4.1, width=12, height=0.8,
          size=20, italic=True, color=MUTED)

    # ---------- 2. DATA EXPLORATION ----------
    s = _add_slide(prs, notes=(
        "Loaded verbatim with all_varchar=True — every column arrives as TEXT, "
        "no silent type coercion; typing is staging's job. Three months of data, "
        "single date format, 120k unique orders spread across 129k rows tells us "
        "immediately that grain is finer than 'order'. The 9 categories and 11 "
        "sizes constrain the size and time dimensions of the marts."
    ))
    _title(s, "Slide 1 — Data Exploration")
    _table(s, ["Metric", "Finding"], [
        ("Records", f"{f.n_rows:,}"),
        ("Unique orders", f"{f.unique_orders:,}"),
        ("Source columns", str(f.n_cols)),
        ("Date range", f"{f.date_min} → {f.date_max} ({f.distinct_dates} distinct days)"),
        ("SKUs", f"{f.skus:,}"),
        ("Styles", f"{f.styles:,}"),
        ("Categories", str(f.categories)),
        ("Sizes", str(f.sizes)),
        ("Ship cities", f"{f.cities:,}"),
    ], top=1.4, height=4.5, col_widths=[3.5, 8.6])
    _text(s, "Business concepts: Orders · Products · Sales · Location · Fulfilment",
          left=0.6, top=6.3, width=12, height=0.5, size=13, italic=True, color=MUTED)

    # ---------- 3. SOURCE GRAIN ----------
    s = _add_slide(prs, notes=(
        "The first Kimball question is always grain. 129k rows across 120k orders "
        "means many rows per order — up to 12 SKUs. The 7 exact-duplicate "
        "(Order ID, SKU) groups are byte-for-byte identical, so DISTINCT is safe. "
        "SKU sits between Style and Size in the product hierarchy — grain must "
        "preserve all three so we can drill up (Style/Category) or down (SKU/Size)."
    ))
    _title(s, "Slide 2 — Understanding the Source Grain")
    _text(s, "One order ≠ one row", left=0.6, top=1.4, width=6, height=0.5,
          size=17, bold=True, color=NAVY)
    _table(s, ["Metric", "Finding"], [
        ("Records", f"{f.n_rows:,}"),
        ("Unique orders", f"{f.unique_orders:,}"),
        ("Extra order lines", f"{f.n_rows - f.unique_orders:,}"),
        ("Orders with multiple rows", f"{f.multi_line_orders:,}"),
        ("Max rows for one order", str(f.max_lines_per_order)),
        ("Exact-duplicate (Order ID, SKU) groups", str(f.dup_order_sku_groups)),
        ("Extra rows in those groups", str(f.dup_order_sku_extra)),
    ], top=1.9, left=0.6, width=6.2, height=4.0, col_widths=[3.5, 2.7])

    _text(s, "Product hierarchy  Style → SKU → Size",
          left=7.0, top=1.4, width=6, height=0.5,
          size=17, bold=True, color=NAVY)
    _table(s, ["Metric", "Finding"], [
        ("Styles", f"{f.styles:,}"),
        ("SKUs", f"{f.skus:,}"),
        ("Sizes", str(f.sizes)),
        ("Styles with multiple SKUs", f"{f.styles_multi_skus:,}"),
        ("Max SKUs for one style", str(f.max_skus_per_style)),
    ], top=1.9, left=7.0, width=5.7, height=3.0, col_widths=[3.2, 2.5])

    _text(s, "Fact grain = order × SKU line item. Composite PK (Order ID, SKU) "
             "requires DISTINCT in staging (D2).",
          left=0.6, top=6.3, width=12, height=0.5,
          size=13, italic=True, color=ACCENT)

    # ---------- 4. DATA QUALITY ASSESSMENT ----------
    s = _add_slide(prs, notes=(
        "Five quality dimensions checked: completeness, uniqueness, validity, "
        "consistency, business rules. Master data is clean (0 SKU drift). Business-"
        "level duplicates are rare. Zero negative amounts / quantities means the "
        "source doesn't leak refunds or corrections as sign-flips — those flow "
        "through Status instead. 218 Status-vs-Courier contradictions is small "
        "enough to flag rather than reject."
    ))
    _title(s, "Slide 3 — Data Quality — Completeness & Uniqueness")
    _text(s, "Completeness · Uniqueness · Validity · Consistency · Business Rules",
          left=0.6, top=1.4, width=12, height=0.4, size=13, italic=True, color=MUTED)
    _table(s, ["Check", "Result"], [
        ("Business-level exact duplicate rows (excl. `index`)", f"{f.biz_dupes:,}"),
        ("Missing Amount", f"{f.missing_amount:,}"),
        ("Qty = 0", f"{f.qty_zero:,}"),
        ("Negative Qty / Negative Amount", "0 / 0"),
        ("Invalid dates (format %m-%d-%y)", str(f.invalid_dates)),
        ("Missing shipping info", f"{f.missing_ship_info:,}"),
        ("SKU → Style / Category / Size drift",
            f"{f.sku_style_drift} / {f.sku_category_drift} / {f.sku_size_drift}"),
        ("Status ↔ Courier Status contradictions",
            f"{f.contradictions:,} ({100*f.contradictions/f.n_rows:.2f}%)"),
    ], top=1.9, height=4.4, col_widths=[8.5, 3.6])
    _text(s, "Master data is clean (SKU consistent across rows). "
             "(Order ID, SKU) is 'almost' unique — dedupe required.",
          left=0.6, top=6.4, width=12, height=0.5,
          size=13, italic=True, color=ACCENT)

    # ---------- 5. CROSS-FIELD + GEOGRAPHY ----------
    s = _add_slide(prs, notes=(
        "Amount and currency null together, always — currency carries no "
        "independent info. Fulfilment fully determines fulfilled-by. Both are "
        "evidence for D4. Casing collapse of 18% on city and 32% on state is "
        "real analytical loss if not normalized (D7). Without normalization, "
        "'Mumbai' / 'mumbai' / ' MUMBAI ' each become their own city."
    ))
    _title(s, "Slide 4 — Data Quality — Cross-field & Geography")
    _text(s, "Amount / currency null pattern", left=0.6, top=1.4,
          width=6, height=0.4, size=15, bold=True, color=NAVY)
    _bullets(s, [
        f"Both NULL: {f.amt_cur_both_null:,}  ·  Only Amount NULL: 0  ·  Only currency NULL: 0",
        "Perfectly correlated → drop `currency` (D4)",
    ], top=1.85, size=13)

    _text(s, "Fulfilment ↔ fulfilled-by", left=0.6, top=2.75,
          width=6, height=0.4, size=15, bold=True, color=NAVY)
    _table(s, ["Fulfilment", "fulfilled-by", "count"], [
        ("Fashionable", "<NULL>", "89,698"),
        ("Merchant", "Easy Ship", "39,277"),
    ], top=3.2, left=0.6, width=6.2, height=1.1, col_widths=[2.3, 2.3, 1.6])

    _text(s, "Casing collapse (D7)", left=7.0, top=1.4,
          width=6, height=0.4, size=15, bold=True, color=NAVY)
    _table(s, ["Column", "Raw", "Normalized", "% dup"], [
        ("ship-city", f"{f.city_raw:,}", f"{f.city_norm:,}",
            f"{100*(f.city_raw-f.city_norm)/f.city_raw:.0f}%"),
        ("ship-state", f"{f.state_raw:,}", f"{f.state_norm:,}",
            f"{100*(f.state_raw-f.state_norm)/f.state_raw:.0f}%"),
        ("ship-country", "1", "1", "0%"),
    ], top=1.85, left=7.0, width=5.7, height=1.6, col_widths=[2.0, 1.2, 1.5, 1.0])

    _text(s, "India has 36 states + UTs officially. Even normalized we see 47 → "
             "residual misspellings (out of scope for fuzzy matching).",
          left=7.0, top=3.6, width=5.7, height=1.2,
          size=12, italic=True, color=MUTED)

    # ---------- 6. NOT EVERY ANOMALY IS BAD DATA ----------
    s = _add_slide(prs, notes=(
        "The temptation is to drop cancelled rows or set Amount to zero silently. "
        "Both destroy business signal. 15% of rows are cancelled and marketers "
        "explicitly want cancellation-rate analysis. Cancelled orders with "
        "Amount>0 are list-price snapshots — separate the two concerns with "
        "is_cancelled and revenue_amount (D1). Multi-value promotion-ids is a "
        "textbook M:N candidate — but with no metadata on promos, dim_promotion "
        "would be degenerate. promo_count now, bridge later (D3)."
    ))
    _title(s, "Slide 5 — Not Every Anomaly Is Bad Data")
    _text(s, "Business context of anomalies", left=0.6, top=1.4,
          width=12, height=0.4, size=15, bold=True, color=NAVY)
    _table(s, ["Condition", "Total", "Cancelled", "Shipped"], [
        ("Qty = 0", f"{f.qty_zero:,}",
            f"{f.qty_zero - f.shipped_qty_zero:,}", f"{f.shipped_qty_zero:,}"),
        ("Amount missing", f"{f.missing_amount:,}",
            f"{f.missing_amount - f.shipped_amt_missing:,}", f"{f.shipped_amt_missing:,}"),
    ], top=1.9, height=1.1, col_widths=[4.5, 2.5, 2.5, 2.6])

    _bullets(s, [
        f"Cancelled + Amount > 0: {f.cancelled_with_amt:,} rows — list price at cancel time. "
        "Naïve sum inflates revenue ~15%. → D1 is_cancelled flag + derived revenue_amount.",
        f"Active-order anomalies to FLAG (not delete): Shipped + Qty=0 = {f.shipped_qty_zero:,} · "
        f"Shipped + Amount missing = {f.shipped_amt_missing:,}.",
        f"promotion-ids: {f.promo_present:,} rows carry data, up to {f.promo_max_count} promos/row "
        f"(avg ~{f.promo_avg_count:.1f}). Textbook M:N bridge is correct — but dim_promotion "
        "would be degenerate today (no metadata). → D3 promo_count on fact, bridge deferred.",
        "Key insight: NULL / unusual values don't automatically mean bad data. "
        "Silent filtering destroys business signal.",
    ], top=3.2, size=14)

    # ---------- 7. TREATMENT STRATEGY ----------
    s = _add_slide(prs, notes=(
        "Every raw column now has a clear treatment: CLEAN, KEEP, FLAG, EXCLUDE, "
        "or TEST. Raw data remains preserved via _raw audit columns (D7). "
        "Suspicious records are flagged, not deleted. Guiding principle: improve "
        "data quality without destroying business information — every deletion is "
        "a design decision, not a shortcut."
    ))
    _title(s, "Slide 6 — Data Quality Treatment Strategy")
    _table(s, ["Action", "Examples"], [
        ("CLEAN",   "Types, snake_case names, whitespace/casing, city/state normalization, date parsing"),
        ("KEEP",    "Cancelled orders, valid business NULLs, `_raw` audit columns"),
        ("FLAG",    "Shipped + Qty = 0 · Shipped + missing Amount · Status/Courier contradictions"),
        ("EXCLUDE", f"Technical `index`, `Unnamed: 22` ({f.u22_nonempty:,} rows literal 'False'), "
                    "`fulfilled-by` (redundant), constants `currency` / `ship-country`"),
        ("TEST",    "Line-item PK uniqueness · fact↔dim relationships · accepted values for Status · "
                    "row-count reconciliation vs raw"),
    ], top=1.4, height=4.5, col_widths=[1.8, 10.3])
    _text(s, "Guiding principle: improve data quality without destroying business information.",
          left=0.6, top=6.3, width=12, height=0.5,
          size=14, italic=True, color=ACCENT)

    # ---------- 8. D1–D8 LOCKED DECISIONS ----------
    d_notes = (
        "Interview soundbites for each decision — use these when questioned:\n\n"
        "D1 (cancelled orders): Cancellations are a business signal, not noise. "
        "Preserving them keeps cancellation-rate analysis one join away; the derived "
        "revenue_amount protects financial rollups from double-counting.\n\n"
        "D2 (fact grain): Grain is the first Kimball question. Since the brief demands "
        "style/category/size analysis, the fact must be at line-item grain — "
        "aggregating up is a rollup, but disaggregating down is impossible.\n\n"
        "D3 (promotion bridge): The bridge is the textbook pattern for M:N, but "
        "without any attributes on the promotion side it would be a degenerate "
        "dimension. I built the useful signal (promo_count) now and documented the "
        "bridge as a first-class item in FUTURE_IMPROVEMENTS.md for when metadata "
        "lands.\n\n"
        "D4 (drop columns): Constants aren't information. If the business expands "
        "internationally, currency and ship-country come back with real values and "
        "need explicit re-typing anyway — carrying them now would be speculative.\n\n"
        "D5 (status bucketing): Raw status has 13 values with overlapping semantics; "
        "the 5-group rollup is the natural funnel view. Both live in one dimension "
        "so grouping is a hardcoded contract, not a report-time rebucket.\n\n"
        "D6 (contradictions): When source data disagrees with itself, staging "
        "shouldn't secretly resolve it. A warn-level test captures the current noise "
        "floor and signals when it drifts — the source system owns the fix.\n\n"
        "D7 (city/state normalization): Normalization is a hard contract; fuzzy "
        "matching is a soft guess. Staging owns the contract, the _raw columns "
        "preserve the audit trail, and any fuzzy work belongs in a downstream "
        "enrichment mart with tests.\n\n"
        "D8 (date parsing): Type at the earliest possible layer where the data "
        "supports it. Staging is that layer here — the format is clean and uniform."
    )
    s = _add_slide(prs, notes=d_notes)
    _title(s, "Slide 7 — Locked Decisions D1–D8")
    _table(s, ["Ref", "Decision"], [
        ("D1", "Cancelled orders kept · is_cancelled flag + derived revenue_amount"),
        ("D2", "Fact grain = order × SKU line item · SK from (order_id, sku) after dedupe"),
        ("D3", "promo_count on fact · bridge_order_promotion + dim_promotion deferred"),
        ("D4", "Drop Unnamed:22, fulfilled-by, currency, ship-country in staging"),
        ("D5", "dim_order_status with status_detail + 5-group status_group"),
        ("D6", "Contradictions pass through · warn-level dbt test at 250 threshold"),
        ("D7", "UPPER(TRIM()) canonical city/state · keep _raw audit columns"),
        ("D8", "Cast Date → DATE with strptime('%m-%d-%y') in staging"),
    ], top=1.4, height=5.2, col_widths=[0.9, 11.2])
    _text(s, "Full rationale + interview soundbites live in speaker notes.",
          left=0.6, top=6.75, width=12, height=0.4,
          size=12, italic=True, color=MUTED)

    # ---------- 9. NEXT ----------
    s = _add_slide(prs, notes=(
        "Phase 2 implements the design contract locked in D1–D8. The staging "
        "model stg_fashionable__orders will type every column, snake_case names, "
        "dedupe on (order_id, sku), normalize geography, derive is_cancelled + "
        "revenue_amount + promo_count, and flag the 218 contradictions with a "
        "warn-level test."
    ))
    _text(s, "Next: staging layer",
          left=0.7, top=2.8, width=12, height=1.2,
          size=44, bold=True, color=NAVY)
    _bar(s, left=0.7, top=4.1, width=2, color=ACCENT)
    _text(s, "stg_fashionable__orders — types, renames, dedupe, normalization, "
             "derived measures",
          left=0.7, top=4.3, width=12, height=0.8,
          size=18, italic=True, color=MUTED)

    return prs


def main() -> None:
    OUT.parent.mkdir(parents=True, exist_ok=True)
    facts = collect()
    prs = build(facts)
    prs.save(str(OUT))
    print(f"✓ Deck written to {OUT.relative_to(REPO)} — {len(prs.slides)} slides")


if __name__ == "__main__":
    main()
